#!/usr/bin/env python3
"""Genera PPT tutoriales Inscripciones y Colegiaturas (admin + padres)."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[2]
OUT = Path(__file__).resolve().parent
ASSETS = Path("/home/mario/.cursor/projects/home-mario-Proyectos/assets")

LOGO_IWC = ROOT / "logos" / "logo-winston-churchill.png"
LOGO_EDU = ROOT / "logos" / "logo-winston-educativo.png"
LOGO_FULL = ROOT / "public" / "bauchers" / "logo_full.png"

SHOT = {
    "cierre": ASSETS / "image-1af4db3c-17c3-4001-a728-f8f686071714.png",
    "cierre_pagado": ASSETS / "image-01bbd502-463d-43fa-8989-fc58cead53e9.png",
    "pasos": ASSETS / "image-7c6bf878-ab7d-49ec-b2ae-bded6990dc8c.png",
    "pago": ASSETS / "image-977c7977-6c96-490c-ac67-f23f0c73781b.png",
    "colegiaturas": ASSETS / "image-9c5378f4-7e80-439b-8230-dcd26067dd97.png",
    "completado": ASSETS / "image-8e83361e-af05-4e01-9087-8c5b8f1d49b8.png",
    "coleg_nuevo": ASSETS / "image-11ed1c38-e759-491a-bca7-0b14237ee4d4.png",
}

# Paleta institucional
NAVY = RGBColor(0x0B, 0x1F, 0x3A)
TEAL = RGBColor(0x0D, 0x94, 0x88)
CYAN = RGBColor(0x0E, 0xA5, 0xE9)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
AMBER = RGBColor(0xD9, 0x77, 0x06)
RED = RGBColor(0xDC, 0x26, 0x26)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SLATE = RGBColor(0x33, 0x41, 0x55)
LIGHT = RGBColor(0xF1, 0xF5, 0xF9)
MUTED = RGBColor(0x64, 0x74, 0x8B)


def blank_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def add_rect(slide, left, top, width, height, fill: RGBColor, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    return shape


def add_round(slide, left, top, width, height, fill: RGBColor):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def set_run(run, text, size=18, bold=False, color=SLATE, font="Calibri"):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_text_box(slide, left, top, width, height, paragraphs, align=PP_ALIGN.LEFT):
    """paragraphs: str | (text, color) | (text, size, bold, color)."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    first = True
    for item in paragraphs:
        if isinstance(item, str):
            text, size, bold, color = item, 16, False, SLATE
        elif len(item) == 2:
            text, color = item
            size, bold = 15, False
        else:
            text, size, bold, color = item
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(6)
        run = p.add_run()
        set_run(run, text, size=size, bold=bold, color=color)
    return box


def header_bar(slide, title: str, subtitle: str | None = None):
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.05), NAVY)
    add_rect(slide, Inches(0), Inches(1.05), Inches(13.333), Inches(0.08), TEAL)
    if LOGO_IWC.exists():
        slide.shapes.add_picture(str(LOGO_IWC), Inches(0.25), Inches(0.12), height=Inches(0.82))
    if LOGO_EDU.exists():
        slide.shapes.add_picture(str(LOGO_EDU), Inches(12.2), Inches(0.12), height=Inches(0.82))
    add_text_box(
        slide,
        Inches(1.3),
        Inches(0.18),
        Inches(10.5),
        Inches(0.45),
        [(title, 26, True, WHITE)],
    )
    if subtitle:
        add_text_box(
            slide,
            Inches(1.3),
            Inches(0.58),
            Inches(10.5),
            Inches(0.35),
            [(subtitle, 13, False, RGBColor(0xA5, 0xF3, 0xFC))],
        )


def footer(slide, page: int, total: int, audience: str):
    add_rect(slide, Inches(0), Inches(7.15), Inches(13.333), Inches(0.35), NAVY)
    add_text_box(
        slide,
        Inches(0.35),
        Inches(7.18),
        Inches(9),
        Inches(0.28),
        [(f"Instituto Winston Churchill  ·  {audience}", 10, False, RGBColor(0xCB, 0xD5, 0xE1))],
    )
    add_text_box(
        slide,
        Inches(10.5),
        Inches(7.18),
        Inches(2.5),
        Inches(0.28),
        [(f"{page} / {total}", 10, False, RGBColor(0xCB, 0xD5, 0xE1))],
        align=PP_ALIGN.RIGHT,
    )


def bullets(slide, left, top, width, height, items: list[tuple[str, RGBColor] | str], size=16):
    paras = []
    for it in items:
        if isinstance(it, tuple):
            text, color = it
        else:
            text, color = it, SLATE
        paras.append((f"•  {text}", size, False, color))
    add_text_box(slide, left, top, width, height, paras)


def add_shot(slide, key: str, left, top, width=None, height=None):
    path = SHOT.get(key)
    if not path or not path.exists():
        return
    kwargs = {}
    if width is not None:
        kwargs["width"] = width
    if height is not None:
        kwargs["height"] = height
    slide.shapes.add_picture(str(path), left, top, **kwargs)


def cover(prs, audience: str, title: str, subtitle: str, badge: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), NAVY)
    add_rect(slide, Inches(0), Inches(0), Inches(0.25), Inches(7.5), TEAL)
    if LOGO_IWC.exists():
        slide.shapes.add_picture(str(LOGO_IWC), Inches(0.7), Inches(0.55), height=Inches(1.35))
    if LOGO_EDU.exists():
        slide.shapes.add_picture(str(LOGO_EDU), Inches(11.3), Inches(0.55), height=Inches(1.35))
    if LOGO_FULL.exists():
        slide.shapes.add_picture(str(LOGO_FULL), Inches(4.2), Inches(0.45), width=Inches(4.8))

    add_round(slide, Inches(0.7), Inches(2.3), Inches(3.2), Inches(0.42), TEAL)
    add_text_box(
        slide,
        Inches(0.7),
        Inches(2.34),
        Inches(3.2),
        Inches(0.35),
        [(badge, 12, True, WHITE)],
        align=PP_ALIGN.CENTER,
    )
    add_text_box(
        slide,
        Inches(0.7),
        Inches(2.95),
        Inches(11.5),
        Inches(1.2),
        [(title, 36, True, WHITE)],
    )
    add_text_box(
        slide,
        Inches(0.7),
        Inches(4.2),
        Inches(11),
        Inches(1.0),
        [(subtitle, 18, False, RGBColor(0xBA, 0xE6, 0xFD))],
    )
    add_text_box(
        slide,
        Inches(0.7),
        Inches(6.5),
        Inches(11),
        Inches(0.4),
        [
            (
                f"Portal: servicios-admin.vercel.app/portal-inscripciones  ·  {audience}",
                13,
                False,
                RGBColor(0x94, 0xA3, 0xB8),
            )
        ],
    )


def build_admin() -> Path:
    prs = blank_prs()
    slides_meta: list[str] = []

    def finish(slide, title_footer: str):
        slides_meta.append(title_footer)

    # 1 cover
    cover(
        prs,
        "Personal administrativo",
        "Tutorial interno:\nInscripciones y Colegiaturas",
        "Guía operativa del portal unificado (reinscritos y nuevo ingreso). "
        "Cómo acompañar a las familias y diagnosticar bloqueos.",
        "USO INTERNO",
    )
    finish(prs.slides[-1], "Portada")

    # 2 objetivo
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(s, "1. Para qué sirve este tutorial", "Personal de servicios / control escolar / caja")
    bullets(
        s,
        Inches(0.6),
        Inches(1.4),
        Inches(12),
        Inches(5),
        [
            ("Explicar el flujo real del portal a padres o en mesa de ayuda.", SLATE),
            ("Identificar por qué un alumno “debe” o no ve colegiaturas.", SLATE),
            ("Distinguir ciclo que se cierra (ej. 22 / 2025-2026) del ciclo nuevo (23 / 2026-2027).", TEAL),
            ("Recordar: los números de ciclo avanzan cada temporada (22→23→24…). No hardcodear.", AMBER),
            ("Saber dónde salen PDF y XML de factura (mismo flujo Banorte/SPEI + timbrado InsForge).", CYAN),
        ],
        size=18,
    )
    finish(s, "Objetivo")

    # 3 mapa
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(s, "2. Mapa del proceso (reinscrito)", "Orden obligatorio")
    cards = [
        ("A", "Cierre ciclo\nanterior", "10 u 11 meses\ndel ciclo actual", AMBER),
        ("B", "4 pasos\nreinscripción", "Solicitud · Reglamento\nPago · Recibo", TEAL),
        ("C", "Colegiaturas\nciclo nuevo", "Concepto 00 +\nmensualidades", CYAN),
    ]
    x = 0.7
    for letter, title, desc, color in cards:
        add_round(s, Inches(x), Inches(1.6), Inches(3.6), Inches(4.2), LIGHT)
        add_round(s, Inches(x + 1.2), Inches(1.85), Inches(1.2), Inches(1.2), color)
        add_text_box(
            s,
            Inches(x + 1.2),
            Inches(2.15),
            Inches(1.2),
            Inches(0.7),
            [(letter, 32, True, WHITE)],
            align=PP_ALIGN.CENTER,
        )
        add_text_box(
            s,
            Inches(x + 0.2),
            Inches(3.3),
            Inches(3.2),
            Inches(1.0),
            [(title, 20, True, NAVY)],
            align=PP_ALIGN.CENTER,
        )
        add_text_box(
            s,
            Inches(x + 0.2),
            Inches(4.5),
            Inches(3.2),
            Inches(1.0),
            [(desc, 14, False, MUTED)],
            align=PP_ALIGN.CENTER,
        )
        x += 4.1
    finish(s, "Mapa")

    # 4 cierre
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(s, "3. Cierre del ciclo anterior", "Antes de reinscribirse debe liquidar 10 u 11 meses")
    add_text_box(
        s,
        Inches(0.5),
        Inches(1.3),
        Inches(6.2),
        Inches(2.2),
        [
            ("Qué valida el sistema", 18, True, NAVY),
            ("Plan 10 meses: 00 + sep→jun + material (16).", 14, False, SLATE),
            ("Plan 11 meses: lo anterior + julio (26).", 14, False, SLATE),
            ("Si la matriz muestra Pagado y el aviso sigue, era bug de match de referencia (ya corregido).", 14, False, AMBER),
        ],
    )
    add_shot(s, "cierre", Inches(6.8), Inches(1.35), width=Inches(6.1))
    finish(s, "Cierre")

    # 5 liquidado
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(s, "4. Cuando ya no adeuda", "Se abre el flujo de reinscripción")
    add_shot(s, "cierre_pagado", Inches(0.5), Inches(1.35), width=Inches(7.4))
    bullets(
        s,
        Inches(8.2),
        Inches(1.5),
        Inches(4.6),
        Inches(5),
        [
            ("Todos los conceptos del plan en Pagado.", GREEN),
            ("Desaparece el aviso amarillo de adeudo.", GREEN),
            ("Empiezan los 4 pasos.", TEAL),
            ("Si sigue bloqueado: Actualizar estado / recargar.", AMBER),
        ],
        size=15,
    )
    finish(s, "Sin adeudo")

    # 6 cuatro pasos
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(s, "5. Los 4 pasos de reinscripción", "Disponible ≠ Completado")
    add_shot(s, "pasos", Inches(0.4), Inches(1.3), width=Inches(7.6))
    bullets(
        s,
        Inches(8.2),
        Inches(1.4),
        Inches(4.7),
        Inches(5.2),
        [
            ("01 Solicitud — datos del alumno.", SLATE),
            ("02 Reglamento — abrir/descargar.", SLATE),
            ("03 Pago — ventanilla / Banorte / SPEI.", SLATE),
            ("04 Recibo final — abrirlo 1 vez.", SLATE),
            ("Colegiaturas del ciclo NUEVO solo con los 4 en Completado.", RED),
        ],
        size=14,
    )
    finish(s, "Pasos")

    # 7 pago y factura
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(s, "6. Pago → timbrado → factura", "Mismo flujo que colegiaturas / admisiones")
    add_shot(s, "pago", Inches(0.5), Inches(1.4), width=Inches(12.2))
    add_text_box(
        s,
        Inches(0.5),
        Inches(3.3),
        Inches(12.2),
        Inches(3.2),
        [
            ("Flujo", 18, True, NAVY),
            ("1) El padre paga (efectivo, comercio electrónico Banorte o SPEI).", 15, False, SLATE),
            ("2) Se registra el pago en pago_detalle.", 15, False, SLATE),
            ("3) Se timbra CFDI y se guarda PDF/XML en InsForge Storage.", 15, False, SLATE),
            ("4) En el paso 03 ya no se muestra “Ver comprobante”: salen botones PDF y XML.", 15, False, TEAL),
            ("Conceptos de reinscripción: 11 / 12 / 13 según diferido o pago completo.", 15, False, CYAN),
        ],
    )
    finish(s, "Factura")

    # 8 colegiaturas
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(s, "7. Colegiaturas del ciclo nuevo", "Se desbloquean al completar los 4 pasos")
    add_shot(s, "completado", Inches(0.5), Inches(1.35), width=Inches(12.2))
    add_text_box(
        s,
        Inches(0.5),
        Inches(3.5),
        Inches(12),
        Inches(3),
        [
            ("Si faltan reglamento o recibo (Disponible), la matriz del ciclo nuevo NO debe mostrarse.", RED),
            ("Tras completar: cuota de inicio de curso (00) y mensualidades del destino (ej. 2026-2027).", TEAL),
            ("Cambridge / Winston USA solo si hay precio real cargado en el ciclo (no inventar montos).", AMBER),
            ("Acordeón “Proceso completado”: Expandir / Collapsar para volver a ver los pasos.", MUTED),
        ],
    )
    finish(s, "Colegiaturas")

    # 9 ciclos
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(s, "8. Ciclos escolares (permanente)", "El número cambia cada año")
    add_round(s, Inches(0.6), Inches(1.5), Inches(12), Inches(5.2), LIGHT)
    add_text_box(
        s,
        Inches(1.0),
        Inches(1.8),
        Inches(11),
        Inches(4.6),
        [
            ("Regla de oro", 22, True, NAVY),
            ("Hoy el ciclo de reinscripción destino es 23 (= 2026-2027).", 17, False, SLATE),
            ("El año próximo será 24, luego 25… Siempre proyección origen → origen+1.", 17, False, TEAL),
            ("es_actual en BD = ciclo en curso de colegiaturas “viejas”.", 17, False, SLATE),
            ("Cierre = ciclo de inscripción − 1 (el que se liquida antes de reinscribir).", 17, False, SLATE),
            ("Nunca fijar “el 23” en código o capacitaciones como valor eterno.", 17, True, RED),
        ],
    )
    finish(s, "Ciclos")

    # 10 checklist
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(s, "9. Checklist rápido mesa de ayuda", "Qué preguntar / revisar")
    left = [
        ("¿Ve aviso de liquidar 2025-2026?", AMBER),
        ("¿Plan 10 u 11 meses?", SLATE),
        ("¿Todos los meses en Pagado?", GREEN),
        ("¿Actualizó estado / recargó?", SLATE),
    ]
    right = [
        ("¿Los 4 pasos en Completado?", TEAL),
        ("¿Pago timbrado? (PDF/XML)", CYAN),
        ("¿Precios del ciclo nuevo cargados?", AMBER),
        ("¿URL: portal-inscripciones?", SLATE),
    ]
    add_text_box(s, Inches(0.6), Inches(1.35), Inches(6), Inches(0.4), [("Bloqueo por adeudo", 16, True, AMBER)])
    bullets(s, Inches(0.6), Inches(1.8), Inches(6), Inches(4.5), left, size=16)
    add_text_box(s, Inches(7), Inches(1.35), Inches(5.8), Inches(0.4), [("Bloqueo por pasos / factura", 16, True, TEAL)])
    bullets(s, Inches(7), Inches(1.8), Inches(5.8), Inches(4.5), right, size=16)
    finish(s, "Checklist")

    # 11 cierre
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(s, "10. Resumen para el equipo", "Una frase por bloque")
    bullets(
        s,
        Inches(0.6),
        Inches(1.4),
        Inches(12),
        Inches(5),
        [
            ("Sin liquidar el ciclo que cierra → no hay reinscripción.", RED),
            ("Sin los 4 pasos completos → no hay colegiaturas del ciclo nuevo.", RED),
            ("Pago en línea = mismo camino que colegiaturas: registro + CFDI + PDF/XML.", TEAL),
            ("Capaciten con captura de pantalla y el número de ciclo de esa temporada.", MUTED),
            ("Dudas técnicas: revisar rama desayunos / despliegue Vercel servicios-admin.", CYAN),
        ],
        size=18,
    )
    finish(s, "Resumen")

    total = len(prs.slides)
    for i, slide in enumerate(prs.slides, start=1):
        if i == 1:
            continue
        footer(slide, i, total, "Personal administrativo · Tutorial interno")

    out = OUT / "Tutorial_Inscripciones_Colegiaturas_ADMIN.pptx"
    prs.save(out)
    return out


def build_padres() -> Path:
    prs = blank_prs()

    cover(
        prs,
        "Familias / Padres de familia",
        "Guía para padres:\nInscripciones y Colegiaturas",
        "Cómo reinscribir a tu hijo(a), pagar en línea o ventanilla, "
        "descargar tu factura y ver las colegiaturas del nuevo ciclo.",
        "PARA FAMILIAS",
    )

    # bienvenida
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(s, "Bienvenido(a)", "Instituto Winston Churchill")
    bullets(
        s,
        Inches(0.6),
        Inches(1.4),
        Inches(12),
        Inches(5),
        [
            ("Todo el proceso de reinscripción e inscripción está en un solo portal.", TEAL),
            ("Enlace: https://servicios-admin.vercel.app/portal-inscripciones", CYAN),
            ("Necesitas el usuario del alumno (número de control) y tu contraseña.", SLATE),
            ("Usa una computadora o celular con internet; Chrome o Edge recomendados.", MUTED),
            ("Guarda siempre tus comprobantes y facturas PDF/XML.", GREEN),
        ],
        size=18,
    )

    # paso a paso overview
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(s, "¿Qué vas a hacer?", "Tres etapas claras")
    etapas = [
        ("1", "Poner al día\nel ciclo actual", "Si aún debes colegiaturas\ndel ciclo que termina,\npágalas primero.", AMBER),
        ("2", "Completar\n4 pasos", "Solicitud, reglamento,\npago de reinscripción\ny recibo final.", TEAL),
        ("3", "Colegiaturas\ndel ciclo nuevo", "Cuota de inicio de curso\ny mensualidades\ndel siguiente ciclo.", CYAN),
    ]
    x = 0.7
    for n, t, d, c in etapas:
        add_round(s, Inches(x), Inches(1.55), Inches(3.7), Inches(4.6), LIGHT)
        add_round(s, Inches(x + 1.25), Inches(1.85), Inches(1.2), Inches(1.2), c)
        add_text_box(
            s,
            Inches(x + 1.25),
            Inches(2.15),
            Inches(1.2),
            Inches(0.7),
            [(n, 28, True, WHITE)],
            align=PP_ALIGN.CENTER,
        )
        add_text_box(
            s,
            Inches(x + 0.25),
            Inches(3.3),
            Inches(3.2),
            Inches(0.9),
            [(t, 18, True, NAVY)],
            align=PP_ALIGN.CENTER,
        )
        add_text_box(
            s,
            Inches(x + 0.25),
            Inches(4.4),
            Inches(3.2),
            Inches(1.4),
            [(d, 13, False, MUTED)],
            align=PP_ALIGN.CENTER,
        )
        x += 4.15

    # cierre
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(s, "Etapa 1 — Liquidar el ciclo que termina", "Ejemplo: 2025-2026")
    add_shot(s, "cierre", Inches(0.4), Inches(1.3), width=Inches(7.5))
    bullets(
        s,
        Inches(8.1),
        Inches(1.4),
        Inches(4.8),
        Inches(5),
        [
            ("Si ves un aviso amarillo, aún hay pagos pendientes del ciclo actual.", AMBER),
            ("Paga un concepto a la vez (baucher o pago en línea).", SLATE),
            ("Cuando todo diga Pagado, ya puedes reinscribir.", GREEN),
            ("Pulsa Actualizar si acabas de pagar.", MUTED),
        ],
        size=14,
    )

    # 4 pasos
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(s, "Etapa 2 — Los 4 pasos", "Todos deben quedar en COMPLETADO")
    add_shot(s, "pasos", Inches(0.4), Inches(1.25), width=Inches(12.4))

    s = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(s, "Detalle de cada paso", "Haz clic en los botones grandes de cada tarjeta")
    rows = [
        ("01 Solicitud", "Revisa o actualiza los datos y guarda.", TEAL),
        ("02 Reglamento", "Abre el reglamento y la carta compromiso.", CYAN),
        ("03 Pago", "Paga la reinscripción (ventanilla, tarjeta o SPEI).", GREEN),
        ("04 Recibo final", "Ábrelo al menos una vez para marcarlo completo.", AMBER),
    ]
    y = 1.35
    for title, desc, color in rows:
        add_round(s, Inches(0.6), Inches(y), Inches(12.1), Inches(1.15), LIGHT)
        add_rect(s, Inches(0.6), Inches(y), Inches(0.18), Inches(1.15), color)
        add_text_box(s, Inches(1.0), Inches(y + 0.18), Inches(11), Inches(0.4), [(title, 18, True, NAVY)])
        add_text_box(s, Inches(1.0), Inches(y + 0.58), Inches(11), Inches(0.4), [(desc, 14, False, SLATE)])
        y += 1.3

    # pago formas
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(s, "Cómo pagar la reinscripción", "Elige la forma que te convenga")
    formas = [
        ("Ventanilla", "Imprime el baucher y paga en el banco.", TEAL),
        ("Tarjeta", "Comercio electrónico Banorte (crédito o débito).", CYAN),
        ("SPEI", "Transferencia desde tu banca en línea.", GREEN),
    ]
    x = 0.7
    for t, d, c in formas:
        add_round(s, Inches(x), Inches(1.6), Inches(3.7), Inches(3.0), LIGHT)
        add_text_box(
            s,
            Inches(x + 0.25),
            Inches(1.9),
            Inches(3.2),
            Inches(0.6),
            [(t, 22, True, c)],
            align=PP_ALIGN.CENTER,
        )
        add_text_box(
            s,
            Inches(x + 0.25),
            Inches(2.7),
            Inches(3.2),
            Inches(1.4),
            [(d, 15, False, SLATE)],
            align=PP_ALIGN.CENTER,
        )
        x += 4.15
    add_text_box(
        s,
        Inches(0.7),
        Inches(5.0),
        Inches(12),
        Inches(1.4),
        [
            (
                "Al pagar con tarjeta o SPEI, el sistema registra el pago, genera la factura electrónica "
                "y te muestra botones PDF y XML para descargarla.",
                16,
                False,
                NAVY,
            )
        ],
    )

    # factura
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(s, "Tu factura PDF y XML", "Cuando el pago ya está registrado")
    add_shot(s, "pago", Inches(0.5), Inches(1.4), width=Inches(12.2))
    add_text_box(
        s,
        Inches(0.5),
        Inches(3.5),
        Inches(12.2),
        Inches(3),
        [
            ("En el paso “Pago de reinscripción” verás botones PDF (rojo) y XML (azul).", TEAL),
            ("Guárdalos en tu correo o carpeta; son tu factura CFDI.", SLATE),
            ("Si aún no aparecen, espera unos minutos y pulsa Actualizar.", AMBER),
        ],
    )

    # recibo y colegiaturas
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(s, "Recibo final y colegiaturas nuevas", "Último paso antes de ver el ciclo nuevo")
    add_shot(s, "completado", Inches(0.5), Inches(1.35), width=Inches(12.2))
    add_text_box(
        s,
        Inches(0.5),
        Inches(3.5),
        Inches(12),
        Inches(3),
        [
            ("Abre el recibo final (paso 04) aunque sea una vez.", AMBER),
            ("Cuando los 4 pasos digan Completado, verás “Proceso de reinscripción completado”.", GREEN),
            ("Ahí puedes Expandir o Collapsar para volver a ver los pasos.", MUTED),
            ("Debajo aparecerán las colegiaturas del ciclo nuevo (cuota de inicio de curso, etc.).", TEAL),
        ],
    )

    # tips
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(s, "Consejos útiles", "Para que el proceso avance sin trabas")
    tips = [
        ("Completa un paso a la vez y espera a ver la etiqueta COMPLETADO.", TEAL),
        ("Si pagaste y no se refleja, usa Actualizar o recarga la página.", AMBER),
        ("No cierres la ventana del banco hasta ver confirmación.", RED),
        ("Conserva PDF y XML de cada pago importante.", GREEN),
        ("¿Dudas? Contacta a servicios escolares o caja de tu plantel.", CYAN),
    ]
    bullets(s, Inches(0.6), Inches(1.4), Inches(12), Inches(5), tips, size=18)

    # cierre familias
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(7.5), NAVY)
    add_rect(s, Inches(0), Inches(0), Inches(0.25), Inches(7.5), TEAL)
    if LOGO_IWC.exists():
        s.shapes.add_picture(str(LOGO_IWC), Inches(5.7), Inches(0.8), height=Inches(1.6))
    add_text_box(
        s,
        Inches(1),
        Inches(2.8),
        Inches(11.3),
        Inches(1.2),
        [("¡Gracias por confiar en Winston!", 32, True, WHITE)],
        align=PP_ALIGN.CENTER,
    )
    add_text_box(
        s,
        Inches(1.5),
        Inches(4.2),
        Inches(10.3),
        Inches(1.2),
        [
            (
                "Portal de inscripciones y colegiaturas\nhttps://servicios-admin.vercel.app/portal-inscripciones",
                16,
                False,
                RGBColor(0xBA, 0xE6, 0xFD),
            )
        ],
        align=PP_ALIGN.CENTER,
    )

    total = len(prs.slides)
    for i, slide in enumerate(prs.slides, start=1):
        if i in (1, total):
            continue
        footer(slide, i, total, "Familias · Guía de inscripciones y colegiaturas")

    out = OUT / "Tutorial_Inscripciones_Colegiaturas_PADRES.pptx"
    prs.save(out)
    return out


def main():
    admin = build_admin()
    padres = build_padres()
    print(admin)
    print(padres)


if __name__ == "__main__":
    main()
